#!/usr/bin/env python3
"""
QBase Strategy Report Generator
Takes screenshots of OOS HTML reports and builds an executive PPT.
"""

import os
import re
import time
from pathlib import Path

from selenium import webdriver
from selenium.webdriver.chrome.options import Options

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ────────────────────────────────────────────────────────────────

BASE_DIR = Path("/Users/simon/Desktop/QBase_v2")
SCREENSHOT_DIR = BASE_DIR / "reports" / "screenshots"
SCREENSHOT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_PPT = Path("/Users/simon/Desktop/QBase_Strategy_Report.pptx")

STRATEGIES = [
    {
        "key": "ag_long",
        "html": BASE_DIR / "research/strong_trend/long/AG/2h/v4_+117.13%/oos.html",
        "name": "AG 做多",
        "quadrant": "白银做多",
        "instrument": "白银(AG)",
        "direction": "做多",
        "freq": "2H",
        "core_indicators": "KAMA + TSI + CMF",
        "signal": "KAMA上升 + TSI为正 + CMF为正 → 三重确认做多",
        "indicators": [
            ("[趋势]", "KAMA(55)", "自适应均线 — 趋势中快速跟随，震荡中自动变慢"),
            ("[动量]", "TSI(35,18)", "真实强度指标 — 双重平滑动量，TSI>0=上涨动力"),
            ("[成交量]", "CMF(50)", "资金流量 — 成交量加权资金方向，正值=资金流入"),
        ],
        "oos_periods": [
            ("2021-12-10", "2022-03-09"),
            ("2022-07-15", "2022-12-14"),
            ("2023-03-10", "2024-05-29"),
        ],
    },
    {
        "key": "ag_short",
        "html": BASE_DIR / "research/strong_trend/short/AG/1h/v18_+44.69%/oos.html",
        "name": "AG 做空",
        "quadrant": "白银做空",
        "instrument": "白银(AG)",
        "direction": "做空",
        "freq": "1H",
        "core_indicators": "HMA + 随机指标",
        "signal": "HMA下行 + 随机指标处于弱势区 → 做空",
        "indicators": [
            ("[趋势]", "HMA(16)", "零延迟均线 — 斜率为负表示下跌趋势"),
            ("[动量]", "随机指标(6,3)", "超买超卖指标 — K值<25表示持续弱势"),
        ],
        "oos_periods": [
            ("2021-05-18", "2021-12-10"),
            ("2022-03-09", "2022-07-15"),
            ("2022-12-14", "2023-03-10"),
        ],
    },
    {
        "key": "i_long",
        "html": BASE_DIR / "research/mild_trend/long/I/4h/v27_+44.71%/oos.html",
        "name": "铁矿 做多",
        "quadrant": "铁矿做多",
        "instrument": "铁矿石(I)",
        "direction": "做多",
        "freq": "4H",
        "core_indicators": "McGinley + 力量指数",
        "signal": "价格在McGinley上方 + 力量指数为正 → 做多",
        "indicators": [
            ("[趋势]", "McGinley均线(20)", "自动调速均线 — 紧密跟随价格，减少假信号"),
            ("[成交量]", "力量指数(13)", "价格变化×成交量 — 正值=买方主导"),
        ],
        "oos_periods": [
            ("2021-11-18", "2022-06-02"),
            ("2022-07-15", "2023-03-13"),
        ],
    },
    {
        "key": "i_short",
        "html": BASE_DIR / "research/mild_trend/short/I/2h/v3_+21.19%/oos.html",
        "name": "铁矿 做空",
        "quadrant": "铁矿做空",
        "instrument": "铁矿石(I)",
        "direction": "做空",
        "freq": "2H",
        "core_indicators": "HMA + Schaff + CMF",
        "signal": "HMA下行 + Schaff动量衰竭 + 资金外流 → 做空",
        "indicators": [
            ("[趋势]", "HMA(40)", "零延迟均线 — 下行表示空头趋势"),
            ("[动量]", "Schaff趋势(50,30)", "综合动量 — 低于75且下降=动量衰竭"),
            ("[成交量]", "资金流量(45)", "CMF资金流向 — 负值=资金持续流出"),
        ],
        "oos_periods": [
            ("2022-06-02", "2022-07-15"),
            ("2023-03-13", "2023-05-25"),
        ],
    },
]

# ── Colors ───────────────────────────────────────────────────────────────

BG_COLOR = RGBColor(0x16, 0x21, 0x3E)
CARD_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_MUTED = RGBColor(0x88, 0x92, 0xB0)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xD4, 0xAA)
RED = RGBColor(0xE9, 0x45, 0x60)
BORDER_COLOR = RGBColor(0x0F, 0x34, 0x60)
TAG_BG = RGBColor(0x0F, 0x34, 0x60)


# ── Step 1: Screenshots ─────────────────────────────────────────────────

def take_screenshots():
    """Take metrics + equity screenshots for each strategy."""
    print("Taking screenshots...")
    opts = Options()
    opts.add_argument("--headless=new")
    opts.add_argument("--window-size=1400,900")
    opts.add_argument("--no-sandbox")
    opts.add_argument("--force-device-scale-factor=2")
    opts.binary_location = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
    driver = webdriver.Chrome(options=opts)

    for s in STRATEGIES:
        html_path = s["html"]
        file_url = f"file://{html_path}"
        key = s["key"]

        # Metrics screenshot
        driver.get(file_url)
        time.sleep(2)
        driver.execute_script("window.scrollTo(0, 0)")
        time.sleep(0.5)
        metrics_path = SCREENSHOT_DIR / f"{key}_metrics.png"
        driver.save_screenshot(str(metrics_path))
        print(f"  Saved {metrics_path.name}")

        # Equity curve screenshot
        driver.execute_script("window.scrollTo(0, 680)")
        time.sleep(0.5)
        equity_path = SCREENSHOT_DIR / f"{key}_equity.png"
        driver.save_screenshot(str(equity_path))
        print(f"  Saved {equity_path.name}")

    driver.quit()
    print("Screenshots done.\n")


# ── Extract metrics from HTML ────────────────────────────────────────────

def extract_metrics(html_path):
    """Extract key metrics from an OOS HTML report."""
    html = open(html_path, encoding="utf-8").read()

    def _search(label_pattern):
        m = re.search(
            label_pattern + r'.*?<div[^>]*class="value[^"]*"[^>]*>([-+]?\d+\.?\d*)',
            html, re.DOTALL,
        )
        return m.group(1) if m else "N/A"

    return {
        "sharpe": _search(r"夏普比率"),
        "total_return": _search(r"总收益"),
        "max_dd": _search(r"最大回撤</div>"),
        "win_rate": _search(r"胜率"),
        "calmar": _search(r"卡玛比率"),
    }


# ── PPT Helpers ──────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.02
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_image_fit(slide, img_path, left, top, width, height):
    """Add image scaled to fit within the box, maintaining aspect ratio."""
    from PIL import Image
    img = Image.open(img_path)
    img_w, img_h = img.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        # Width-constrained
        actual_w = width
        actual_h = int(width / img_ratio)
    else:
        actual_h = height
        actual_w = int(height * img_ratio)
    slide.shapes.add_picture(str(img_path), left, top, actual_w, actual_h)


# ── Build PPT ────────────────────────────────────────────────────────────

def build_ppt():
    print("Building PPT...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank

    # ── Slide 1: Cover ───────────────────────────────────────────────────

    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)

    # Accent line
    add_rect(slide, Inches(3.5), Inches(2.6), Inches(3), Pt(3), GOLD)

    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(1),
                 "QBase 量化策略研究报告", font_size=36, color=TEXT_WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.6),
                 "四象限最优策略 \u00b7 OOS 样本外验证", font_size=18,
                 color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.5),
                 "AlphaForge V7.2 \u00b7 2026年4月", font_size=13,
                 color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary ───────────────────────────────────────

    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)

    # Title
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(5), Inches(0.6),
                 "策略总览", font_size=24, color=GOLD, bold=True)

    # Accent line under title
    add_rect(slide, Inches(0.6), Inches(0.85), Inches(1.2), Pt(3), GOLD)

    # Table background card
    card = add_rect(slide, Inches(0.4), Inches(1.2), Inches(9.2), Inches(4.4), CARD_COLOR, BORDER_COLOR)

    # Table header
    headers = ["象限", "品种", "方向", "频率", "Sharpe", "收益率", "核心指标"]
    col_widths = [0.8, 0.8, 0.6, 0.6, 0.8, 0.9, 4.2]  # in inches, total ~8.7
    col_starts = [0.65]
    for w in col_widths[:-1]:
        col_starts.append(col_starts[-1] + w)

    header_top = Inches(1.35)
    for i, h in enumerate(headers):
        add_text_box(slide, Inches(col_starts[i]), header_top, Inches(col_widths[i]), Inches(0.4),
                     h, font_size=12, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Header separator
    add_rect(slide, Inches(0.6), Inches(1.75), Inches(8.8), Pt(1), BORDER_COLOR)

    # Build rows dynamically from STRATEGIES + extracted metrics
    for row_idx, s in enumerate(STRATEGIES):
        m = extract_metrics(s["html"])
        row_top = Inches(1.9 + row_idx * 0.85)
        # Row background (alternating)
        if row_idx % 2 == 0:
            add_rect(slide, Inches(0.5), row_top - Inches(0.05),
                     Inches(9.0), Inches(0.75), RGBColor(0x14, 0x1E, 0x38))

        row = (
            s["quadrant"], s["instrument"], s["direction"], s["freq"],
            m["sharpe"], f'+{m["total_return"]}%', s["core_indicators"],
        )
        for col_idx, val in enumerate(row):
            color = TEXT_WHITE
            if col_idx in (4, 5):
                color = GREEN
            fs = 13 if col_idx < 6 else 11
            add_text_box(slide, Inches(col_starts[col_idx]), row_top, Inches(col_widths[col_idx]),
                         Inches(0.5), val, font_size=fs, color=color,
                         alignment=PP_ALIGN.CENTER if col_idx < 6 else PP_ALIGN.LEFT)

    # ── OOS Time Period Summary ──────────────────────────────────────────
    oos_card_top = Inches(5.6)
    add_rect(slide, Inches(0.4), oos_card_top, Inches(9.2), Inches(1.6), CARD_COLOR, BORDER_COLOR)

    add_text_box(slide, Inches(0.6), oos_card_top + Inches(0.05),
                 Inches(3), Inches(0.35),
                 "OOS 样本外验证时段", font_size=13, color=GOLD, bold=True)

    for oos_idx, s in enumerate(STRATEGIES):
        periods = s["oos_periods"]
        periods_str = "  |  ".join(f"{p[0]} ~ {p[1]}" for p in periods)
        line = f'{s["quadrant"]}（{len(periods)}段）：{periods_str}'
        add_text_box(slide, Inches(0.7), oos_card_top + Inches(0.4 + oos_idx * 0.28),
                     Inches(8.8), Inches(0.28),
                     line, font_size=10, color=TEXT_LIGHT)

    # Footer note
    add_text_box(slide, Inches(0.6), Inches(7.25), Inches(8.8), Inches(0.3),
                 "* 所有数据基于样本外（OOS）回测验证，使用 AlphaForge V7.2 工业级模式（含真实交易成本、保证金、滑点）",
                 font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # ── Slides 3-6: Strategy detail slides ───────────────────────────────

    for s in STRATEGIES:
        metrics = extract_metrics(s["html"])
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, BG_COLOR)

        # Title
        title_text = f'{s["name"]}  ({s["freq"]})'
        add_text_box(slide, Inches(0.4), Inches(0.2), Inches(6), Inches(0.5),
                     title_text, font_size=24, color=GOLD, bold=True)
        add_rect(slide, Inches(0.4), Inches(0.7), Inches(1.5), Pt(3), GOLD)

        # ── Top left: Metrics screenshot (60% width) ──
        metrics_img = SCREENSHOT_DIR / f'{s["key"]}_metrics.png'
        img_left = Inches(0.3)
        img_top = Inches(0.9)
        img_w = Inches(5.7)
        img_h = Inches(2.8)
        if metrics_img.exists():
            add_image_fit(slide, metrics_img, img_left, img_top, img_w, img_h)

        # ── Top right: Strategy info card (40% width) ──
        card_left = Inches(6.1)
        card_top = Inches(0.9)
        card_w = Inches(3.6)
        card_h = Inches(2.8)
        add_rect(slide, card_left, card_top, card_w, card_h, CARD_COLOR, BORDER_COLOR)

        # Strategy name in card
        add_text_box(slide, card_left + Inches(0.2), card_top + Inches(0.1),
                     card_w - Inches(0.4), Inches(0.4),
                     s["name"], font_size=16, color=TEXT_WHITE, bold=True)

        # Metrics in card
        metric_lines = [
            f'频率：{s["freq"]}',
            f'Sharpe：{metrics["sharpe"]}',
            f'收益率：{metrics["total_return"]}%',
            f'最大回撤：{metrics["max_dd"]}%',
            f'胜率：{metrics["win_rate"]}%',
            f'卡玛比率：{metrics["calmar"]}',
        ]
        for mi, line in enumerate(metric_lines):
            y_off = card_top + Inches(0.55 + mi * 0.33)
            add_text_box(slide, card_left + Inches(0.2), y_off,
                         card_w - Inches(0.4), Inches(0.3),
                         line, font_size=12, color=TEXT_LIGHT)

        # Signal logic
        add_rect(slide, card_left + Inches(0.15), card_top + Inches(2.55),
                 card_w - Inches(0.3), Pt(1), BORDER_COLOR)

        # ── Bottom left: Equity curve screenshot (60% width) ──
        equity_img = SCREENSHOT_DIR / f'{s["key"]}_equity.png'
        eq_top = Inches(3.9)
        eq_h = Inches(2.8)
        if equity_img.exists():
            add_image_fit(slide, equity_img, img_left, eq_top, img_w, eq_h)

        # ── Bottom right: Indicator cards (40% width) ──
        ind_top = Inches(3.9)
        ind_h_total = Inches(2.8)
        add_rect(slide, card_left, ind_top, card_w, ind_h_total, CARD_COLOR, BORDER_COLOR)

        # Indicators title
        add_text_box(slide, card_left + Inches(0.15), ind_top + Inches(0.08),
                     card_w - Inches(0.3), Inches(0.35),
                     "信号指标", font_size=14, color=GOLD, bold=True)

        for ii, (tag, name, desc) in enumerate(s["indicators"]):
            y = ind_top + Inches(0.45 + ii * 0.65)

            # Tag pill
            tag_shape = add_rect(slide, card_left + Inches(0.15), y,
                                 Inches(0.65), Inches(0.25), TAG_BG, BORDER_COLOR)
            tag_shape.text_frame.paragraphs[0].text = tag
            tag_shape.text_frame.paragraphs[0].font.size = Pt(9)
            tag_shape.text_frame.paragraphs[0].font.color.rgb = GOLD
            tag_shape.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
            tag_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Indicator name
            add_text_box(slide, card_left + Inches(0.85), y,
                         card_w - Inches(1.0), Inches(0.25),
                         name, font_size=12, color=TEXT_WHITE, bold=True)

            # Indicator description
            add_text_box(slide, card_left + Inches(0.15), y + Inches(0.25),
                         card_w - Inches(0.3), Inches(0.35),
                         desc, font_size=10, color=TEXT_MUTED)

        # Signal summary at bottom of indicator card
        sig_y = ind_top + Inches(0.45 + len(s["indicators"]) * 0.65 + 0.05)
        add_rect(slide, card_left + Inches(0.1), sig_y - Inches(0.05),
                 card_w - Inches(0.2), Pt(1), BORDER_COLOR)
        add_text_box(slide, card_left + Inches(0.15), sig_y + Inches(0.05),
                     card_w - Inches(0.3), Inches(0.4),
                     s["signal"], font_size=10, color=GREEN, bold=True)

        # ── OOS Time Periods ──
        oos_y = Inches(6.85)
        add_rect(slide, Inches(0.3), oos_y, Inches(9.4), Inches(0.5), CARD_COLOR, BORDER_COLOR)
        periods = s["oos_periods"]
        periods_str = "  |  ".join(f"{p[0]} ~ {p[1]}" for p in periods)
        oos_text = f"OOS验证时段（{len(periods)}段）：{periods_str}"
        add_text_box(slide, Inches(0.5), oos_y + Inches(0.08),
                     Inches(9.0), Inches(0.35),
                     oos_text, font_size=11, color=TEXT_LIGHT)

    # ── Save ─────────────────────────────────────────────────────────────

    prs.save(str(OUTPUT_PPT))
    print(f"\nPPT saved to {OUTPUT_PPT}")


# ── Main ─────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    take_screenshots()
    build_ppt()
    print("Done!")
