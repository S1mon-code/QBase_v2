"""Generate strategy presentation PPT for the 4-quadrant best strategies.

Saves to ~/Desktop/QBase_Strategy_Report.pptx
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PROJECT = Path("/Users/simon/Desktop/QBase_v2")

# ── Colors ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x25, 0x25, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xEF, 0x53, 0x50)
GOLD = RGBColor(0xFF, 0xAB, 0x40)
BLUE = RGBColor(0x42, 0xA5, 0xF5)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)

# ── Strategy data ──
STRATEGIES = [
    {
        "quadrant": "AG Long",
        "instrument": "白银 (AG)",
        "direction": "做多",
        "regime": "Strong Trend",
        "freq": "2H",
        "version": "v4",
        "oos_sharpe": 1.271,
        "oos_return": "+129.89%",
        "strategy_name": "long_long_AG_2h_v4",
        "class_name": "StrongTrendLongAG2hV4",
        "indicators": [
            {"name": "KAMA(55)", "category": "趋势", "color": GOLD,
             "desc": "Kaufman 自适应均线 — 根据效率比率自动调节平滑速度，在趋势中快速跟随，在震荡中减少噪音。周期55适配2H白银的中期趋势节奏。"},
            {"name": "TSI(35,18)", "category": "动量", "color": PURPLE,
             "desc": "True Strength Index — 双重 EMA 平滑的动量指标，过滤短期噪音。TSI > 0 表示上升趋势动量，TSI 上穿信号线为买入信号。"},
            {"name": "CMF(50)", "category": "成交量", "color": GREEN,
             "desc": "Chaikin Money Flow — 成交量加权的资金流向指标。CMF > 0 表示资金净流入（买方力量占优），确认趋势有量能支撑。"},
        ],
        "signal_logic": "KAMA 上升（斜率为正）+ TSI > 0 + CMF > 0 → 三重确认做多",
        "pros": [
            "KAMA 自适应特性减少假信号，白银波动性变化大时尤其有效",
            "三指标覆盖趋势+动量+量能三维度，信号质量高",
            "OOS Sharpe 1.271 + Return 129.89%，收益风险比极佳",
        ],
        "cons": [
            "长周期参数（KAMA=55, CMF=50）导致入场较慢，可能错过趋势初期",
            "白银暴跌时平仓延迟较大（chandelier_mult=3.0 止损较宽）",
            "仅在强趋势环境下有效，均值回归行情会产生持续亏损",
        ],
        "research_path": "research/long/long/AG/2h/v4_+129.89%",
    },
    {
        "quadrant": "AG Short",
        "instrument": "白银 (AG)",
        "direction": "做空",
        "regime": "Strong Trend",
        "freq": "1H",
        "version": "v13",
        "oos_sharpe": None,
        "oos_return": "+56.62%",
        "strategy_name": "long_short_AG_1h_v13",
        "class_name": "StrongTrendShortAG1hV13",
        "indicators": [
            {"name": "MACD(6,16,5)", "category": "动量", "color": BLUE,
             "desc": "快速 MACD — 短周期参数(6/16/5)捕捉1H白银的快速动量反转。MACD线 < 信号线表示空头动量占优。"},
            {"name": "Force Index(6)", "category": "成交量", "color": PURPLE,
             "desc": "Force Index — 价格变化 × 成交量，衡量买卖力量的强度。FI < 0 表示卖方力量主导，确认下跌有成交量支撑。"},
        ],
        "signal_logic": "MACD线 < 信号线 + Force Index < 0 → 动量+量能双确认做空",
        "pros": [
            "快速参数适合1H白银的高波动性，捕捉短线做空机会",
            "Force Index 直接量化卖出力量，避免无量下跌的假信号",
            "OOS Return +56.62%，在做空环境中表现优异",
        ],
        "cons": [
            "快速参数容易过度交易，在震荡市中产生频繁止损",
            "白银流动性在夜盘较低，1H信号可能受滑点影响",
            "做空策略天然面临无限风险，依赖 Chandelier Exit 止损",
        ],
        "research_path": "research/long/short/AG/1h/v13_+56.62%",
    },
    {
        "quadrant": "I Long",
        "instrument": "铁矿石 (I)",
        "direction": "做多",
        "regime": "Mild Trend",
        "freq": "1H",
        "version": "v23",
        "oos_sharpe": 0.829,
        "oos_return": "+32.21%",
        "strategy_name": "long_long_I_1h_v23",
        "class_name": "MildTrendLongI1hV23",
        "indicators": [
            {"name": "Coppock Curve(10,14)", "category": "动量", "color": PURPLE,
             "desc": "Coppock 曲线 — 两个 ROC 的加权移动平均，原始设计用于捕捉市场底部反转买入信号。从零线下方转正为强烈买入信号。"},
            {"name": "OI Flow(20)", "category": "持仓量", "color": GREEN,
             "desc": "OI Flow — 以价格方向加权的持仓量流向指标。OI Flow > Signal 表示机构在建仓方向与趋势一致，有基本面资金支撑。"},
        ],
        "signal_logic": "Coppock > 0 且上升 + OI Flow > Signal → 动量反转+机构建仓确认做多",
        "pros": [
            "Coppock Curve 是经典的底部捕捉指标，适合温和趋势中的低位买入",
            "OI Flow 提供期货市场独有的持仓量分析，区分散户和机构行为",
            "信号逻辑分层（强/弱信号），仓位管理灵活",
        ],
        "cons": [
            "Coppock Curve 为滞后指标，趋势已走大半时才发出信号",
            "OI 数据在某些时段可能不准确（换月期间 OI 异常波动）",
            "OOS Return +32.21% 相对 AG Long 偏低，铁矿波动率较低",
        ],
        "research_path": "research/long/long/I/1h/v23_+32.21%",
    },
    {
        "quadrant": "I Short",
        "instrument": "铁矿石 (I)",
        "direction": "做空",
        "regime": "Mild Trend",
        "freq": "2H",
        "version": "v3",
        "oos_sharpe": 1.198,
        "oos_return": "+21.19%",
        "strategy_name": "long_short_I_2h_v3",
        "class_name": "MildTrendShortI2hV3",
        "indicators": [
            {"name": "HMA(40)", "category": "趋势", "color": GOLD,
             "desc": "Hull 移动均线 — 利用加权移动平均的差值消除滞后，比 SMA/EMA 更快响应趋势变化。HMA 下降表示短期趋势转空。"},
            {"name": "Schaff Trend Cycle(50,30)", "category": "动量", "color": PURPLE,
             "desc": "Schaff 趋势周期 — 将 MACD 通过双重随机平滑处理，输出 0-100 区间值。< 75 且下降表示动量衰竭，空头信号增强。"},
            {"name": "CMF(45)", "category": "成交量", "color": GREEN,
             "desc": "Chaikin Money Flow — CMF < 0 表示资金持续流出，卖方主导市场。结合趋势和动量共同确认做空方向。"},
        ],
        "signal_logic": "HMA 下降 + Schaff < 75 且下降 + CMF < 0 → 趋势+动量+量能三重做空确认",
        "pros": [
            "HMA 低延迟特性使做空入场更及时，减少追空风险",
            "三维信号（趋势+动量+量能）过滤效果好，OOS Sharpe 1.198",
            "适合铁矿的温和下跌趋势，信号不会在暴跌中过度加仓",
        ],
        "cons": [
            "HMA 周期40 + CMF 周期45 在2H频率上约等于4-5天，可能错过日内反转",
            "Schaff Trend Cycle 在快速反转时反应不够快（双重平滑的代价）",
            "做空铁矿在供给侧改革等政策驱动的上涨中可能面临较大回撤",
        ],
        "research_path": "research/long/short/I/2h/v3_+21.19%",
    },
]


def set_slide_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_indicator_card(slide, left, top, indicator):
    """Add a styled indicator card."""
    # Card background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(4.2), Inches(1.1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()

    # Category tag
    add_text_box(slide, left + 0.15, top + 0.08, 0.8, 0.3,
                 indicator["category"], font_size=10, color=indicator["color"], bold=True)

    # Name
    add_text_box(slide, left + 1.0, top + 0.05, 3.0, 0.3,
                 indicator["name"], font_size=13, color=WHITE, bold=True)

    # Description
    add_text_box(slide, left + 0.15, top + 0.35, 3.9, 0.75,
                 indicator["desc"], font_size=9, color=GRAY)


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 1.0, 9.0, 1.0,
                 "QBase_v2 策略研报", font_size=36, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 2.0, 9.0, 0.5,
                 "四象限最优策略分析", font_size=20, color=GOLD,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 3.0, 9.0, 0.5,
                 "AG Long  |  AG Short  |  I Long  |  I Short", font_size=16, color=GRAY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 4.5, 9.0, 1.0,
                 "回测引擎：AlphaForge V7.6.1 Industrial Mode\n"
                 "验证体系：6层验证 + 5层归因 + Signal Blending\n"
                 "报告日期：2026-04-03",
                 font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


def create_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.5,
                 "四象限策略总览", font_size=24, color=WHITE, bold=True)

    # Table
    rows, cols = 5, 6
    tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(9.0), Inches(3.0)).table

    headers = ["象限", "品种/方向", "频率", "版本", "OOS Sharpe", "OOS Return"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = GOLD
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    data = [
        ["AG Long", "白银 做多", "2H", "v4", "1.271", "+129.89%"],
        ["AG Short", "白银 做空", "1H", "v13", "—", "+56.62%"],
        ["I Long", "铁矿 做多", "1H", "v23", "0.829", "+32.21%"],
        ["I Short", "铁矿 做空", "2H", "v3", "1.198", "+21.19%"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
                if c == 5:
                    p.font.color.rgb = GREEN if val.startswith("+") else RED

    # Style table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if r > 0 else RGBColor(0x30, 0x30, 0x50)

    add_text_box(slide, 0.5, 4.3, 9.0, 2.0,
                 "选择标准：每个象限中 OOS Sharpe 最高的策略\n"
                 "• AG Long 2h v4：三指标（KAMA+TSI+CMF）覆盖趋势+动量+量能，收益率最高\n"
                 "• AG Short 1h v13：快速 MACD+Force Index 捕捉做空机会，1H 频率最灵活\n"
                 "• I Long 1h v23：Coppock+OI Flow 结合底部反转和机构持仓分析\n"
                 "• I Short 2h v3：HMA+Schaff+CMF 三维做空确认，Sharpe 1.198 最稳健",
                 font_size=10, color=GRAY)


def create_strategy_slides(prs, strategy):
    """Create 2 slides per strategy: overview + indicators."""
    s = strategy

    # ── Slide 1: Strategy Overview ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_text_box(slide, 0.3, 0.2, 5.0, 0.5,
                 f"{s['quadrant']} — {s['instrument']} {s['direction']}",
                 font_size=22, color=WHITE, bold=True)
    add_text_box(slide, 5.5, 0.2, 4.0, 0.5,
                 f"{s['regime']} | {s['freq']} | {s['version']}",
                 font_size=14, color=GOLD, alignment=PP_ALIGN.RIGHT)

    # Metrics cards
    metrics = [
        ("OOS Sharpe", f"{s['oos_sharpe']:.3f}" if s['oos_sharpe'] else "—", BLUE),
        ("OOS Return", s['oos_return'], GREEN if s['oos_return'].startswith("+") else RED),
        ("频率", s['freq'], WHITE),
        ("Regime", s['regime'], GOLD),
    ]

    for i, (label, value, color) in enumerate(metrics):
        left = 0.3 + i * 2.35
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(0.9), Inches(2.15), Inches(0.8),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, 0.92, 1.95, 0.3, label, font_size=10, color=GRAY)
        add_text_box(slide, left + 0.1, 1.2, 1.95, 0.4, value, font_size=18, color=color, bold=True)

    # Signal logic
    add_text_box(slide, 0.3, 1.9, 9.0, 0.3, "信号逻辑", font_size=14, color=GOLD, bold=True)
    add_text_box(slide, 0.3, 2.2, 9.0, 0.5, s['signal_logic'], font_size=12, color=WHITE)

    # Indicators section
    add_text_box(slide, 0.3, 2.8, 9.0, 0.3, "指标组合", font_size=14, color=GOLD, bold=True)

    for i, ind in enumerate(s['indicators']):
        add_indicator_card(slide, 0.3 + (i % 2) * 4.5, 3.2 + (i // 2) * 1.2, ind)

    # Report reference
    add_text_box(slide, 0.3, 6.8, 9.0, 0.3,
                 f"报告路径: {s['research_path']}/oos.html",
                 font_size=9, color=GRAY)

    # ── Slide 2: Pros & Cons ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, BG_DARK)

    add_text_box(slide2, 0.3, 0.2, 9.0, 0.5,
                 f"{s['quadrant']} — 优缺点分析",
                 font_size=22, color=WHITE, bold=True)

    # Pros
    add_text_box(slide2, 0.3, 0.9, 4.5, 0.3, "✓ 优势", font_size=16, color=GREEN, bold=True)
    for i, pro in enumerate(s['pros']):
        shape = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(1.3 + i * 0.9), Inches(4.3), Inches(0.75),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1B, 0x2E, 0x1B)
        shape.line.fill.background()
        add_text_box(slide2, 0.5, 1.35 + i * 0.9, 3.9, 0.7, pro, font_size=10, color=WHITE)

    # Cons
    add_text_box(slide2, 5.0, 0.9, 4.5, 0.3, "✗ 劣势", font_size=16, color=RED, bold=True)
    for i, con in enumerate(s['cons']):
        shape = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.0), Inches(1.3 + i * 0.9), Inches(4.3), Inches(0.75),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2E, 0x1B, 0x1B)
        shape.line.fill.background()
        add_text_box(slide2, 5.2, 1.35 + i * 0.9, 3.9, 0.7, con, font_size=10, color=WHITE)

    # Strategy details
    add_text_box(slide2, 0.3, 4.3, 9.0, 0.3, "策略详情", font_size=14, color=GOLD, bold=True)

    details = (
        f"策略名称: {s['strategy_name']}\n"
        f"类名: {s['class_name']}\n"
        f"信号维度: {', '.join(ind['category'] for ind in s['indicators'])}\n"
        f"信号范围: {'[0, 1] 只做多' if s['direction'] == '做多' else '[-1, 0] 只做空'}\n"
        f"可优化参数: {len(s['indicators']) + 1} 个（含 chandelier_mult）\n"
        f"回测模式: AlphaForge V7.6.1 Industrial（含成交量适应价差、动态保证金、方向不对称冲击）"
    )
    add_text_box(slide2, 0.3, 4.65, 9.0, 2.0, details, font_size=10, color=GRAY)


def create_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 1.0, 9.0, 0.5,
                 "总结与下一步", font_size=28, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    summary = (
        "◆ 系统概览\n"
        "  · 项目：QBase_v2 黑色系+贵金属 多策略系统\n"
        "  · 品种：铁矿石 (I)、白银 (AG)\n"
        "  · 策略总数：390 个（4 象限 × 4 频率 × 多版本）\n"
        "  · 回测引擎：AlphaForge V7.6.1 Industrial Mode\n\n"
        "◆ 四象限最优组合\n"
        "  · AG Long 2h v4: KAMA+TSI+CMF → Sharpe 1.27, Return +129.89%\n"
        "  · AG Short 1h v13: MACD+Force Index → Return +56.62%\n"
        "  · I Long 1h v23: Coppock+OI Flow → Sharpe 0.83, Return +32.21%\n"
        "  · I Short 2h v3: HMA+Schaff+CMF → Sharpe 1.20, Return +21.19%\n\n"
        "◆ 下一步计划\n"
        "  · Portfolio 构建：Signal Blending（Carver 标准） + Regime Allocator\n"
        "  · 扩展品种：RB（螺纹钢）、HC（热卷）、J（焦炭）、JM（焦煤）\n"
        "  · Paper Trading 验证 → 实盘部署"
    )
    add_text_box(slide, 0.8, 1.8, 8.4, 5.0, summary, font_size=13, color=WHITE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_overview_slide(prs)

    for s in STRATEGIES:
        create_strategy_slides(prs, s)

    create_summary_slide(prs)

    output = Path("/Users/simon/Desktop/QBase_Strategy_Report.pptx")
    prs.save(str(output))
    print(f"PPT saved: {output}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
