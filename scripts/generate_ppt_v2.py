"""Generate strategy PPT with OOS report screenshots.

Takes screenshots of metrics cards and equity curve from each OOS HTML report,
then builds PPT with those images + indicator analysis.
"""
from __future__ import annotations

import re
import time
from pathlib import Path
from urllib.parse import quote

from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PROJECT = Path("/Users/simon/Desktop/QBase_v2")
SCREENSHOT_DIR = PROJECT / "reports" / "screenshots"
SCREENSHOT_DIR.mkdir(parents=True, exist_ok=True)

# ── Colors ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x25, 0x25, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xEF, 0x53, 0x50)
GOLD = RGBColor(0xFF, 0xAB, 0x40)
BLUE = RGBColor(0x42, 0xA5, 0xF5)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)


def extract_metrics_from_html(html_path: Path) -> dict:
    """Extract key metrics from OOS HTML report."""
    html = html_path.read_text(encoding="utf-8")
    metrics = {}
    patterns = {
        "sharpe": r'夏普比率.*?<div[^>]*class="value[^"]*"[^>]*>([-+]?\d+\.\d+)',
        "return": r'总收益.*?<div[^>]*class="value[^"]*"[^>]*>([-+]?\d+\.\d+)%',
        "max_dd": r'最大回撤.*?<div[^>]*class="value[^"]*"[^>]*>([-+]?\d+\.\d+)%',
        "calmar": r'卡玛比率.*?<div[^>]*class="value[^"]*"[^>]*>([-+]?\d+\.\d+)',
        "win_rate": r'胜率.*?<div[^>]*class="value[^"]*"[^>]*>([-+]?\d+\.\d+)%',
        "trades": r'交易次数.*?<div[^>]*class="value[^"]*"[^>]*>(\d+)',
        "annualized": r'年化收益.*?<div[^>]*class="value[^"]*"[^>]*>([-+]?\d+\.\d+)%',
        "volatility": r'年化波动率.*?<div[^>]*class="value[^"]*"[^>]*>([-+]?\d+\.\d+)%',
    }
    for key, pattern in patterns.items():
        m = re.search(pattern, html, re.DOTALL)
        if m:
            metrics[key] = m.group(1)
    return metrics


def take_screenshots(driver, strategies):
    """Take metrics + equity curve screenshots for each strategy."""
    for s in strategies:
        oos_html = PROJECT / s["research_path"] / "oos.html"
        if not oos_html.exists():
            print(f"  [skip] {s['quadrant']}: oos.html not found")
            continue

        url = "file://" + str(oos_html).replace("%", "%25").replace("+", "%2B")
        print(f"  Loading {s['quadrant']} {s['version']}...")
        driver.get(url)
        time.sleep(3)

        # 1. Full page screenshot (for metrics cards at top)
        # Set window tall enough to capture metrics
        driver.set_window_size(1400, 900)
        time.sleep(1)

        # Metrics screenshot (top of page)
        metrics_path = SCREENSHOT_DIR / f"{s['strategy_name']}_metrics.png"
        driver.save_screenshot(str(metrics_path))
        print(f"    Metrics: {metrics_path.name}")

        # 2. Scroll down to equity curve and screenshot
        # Equity curve is typically the 2nd section
        driver.set_window_size(1400, 800)
        driver.execute_script("window.scrollTo(0, 650)")
        time.sleep(1)

        equity_path = SCREENSHOT_DIR / f"{s['strategy_name']}_equity.png"
        driver.save_screenshot(str(equity_path))
        print(f"    Equity: {equity_path.name}")

        s["metrics_img"] = metrics_path
        s["equity_img"] = equity_path


STRATEGIES = [
    {
        "quadrant": "AG Long",
        "instrument": "白银 (AG)",
        "direction": "做多",
        "regime": "Strong Trend",
        "freq": "2H",
        "version": "v4",
        "strategy_name": "strong_trend_long_AG_2h_v4",
        "class_name": "StrongTrendLongAG2hV4",
        "research_path": "research/strong_trend/long/AG/2h/v4_+129.89%",
        "indicators": [
            {"name": "KAMA(55)", "category": "趋势", "color": GOLD,
             "desc": "Kaufman 自适应均线 — 根据效率比率自动调节平滑速度，在趋势中快速跟随，在震荡中减少噪音"},
            {"name": "TSI(35,18)", "category": "动量", "color": PURPLE,
             "desc": "True Strength Index — 双重 EMA 平滑的动量指标，TSI > 0 表示上升趋势动量"},
            {"name": "CMF(50)", "category": "成交量", "color": GREEN,
             "desc": "Chaikin Money Flow — 成交量加权的资金流向，CMF > 0 表示资金净流入"},
        ],
        "signal_logic": "KAMA 上升 + TSI > 0 + CMF > 0 → 三重确认做多",
        "pros": ["KAMA 自适应特性减少假信号", "三维度覆盖（趋势+动量+量能）", "OOS Return +129.89%，收益最高"],
        "cons": ["长周期参数入场偏慢", "止损较宽（chandelier=3.0）", "均值回归行情会亏损"],
    },
    {
        "quadrant": "AG Short",
        "instrument": "白银 (AG)",
        "direction": "做空",
        "regime": "Strong Trend",
        "freq": "1H",
        "version": "v13",
        "strategy_name": "strong_trend_short_AG_1h_v13",
        "class_name": "StrongTrendShortAG1hV13",
        "research_path": "research/strong_trend/short/AG/1h/v13_+56.62%",
        "indicators": [
            {"name": "MACD(6,16,5)", "category": "动量", "color": BLUE,
             "desc": "快速 MACD — 短周期参数捕捉1H白银的快速动量反转，MACD线 < 信号线为空头"},
            {"name": "Force Index(6)", "category": "成交量", "color": PURPLE,
             "desc": "Force Index — 价格变化×成交量，FI < 0 表示卖方力量主导"},
        ],
        "signal_logic": "MACD线 < 信号线 + Force Index < 0 → 动量+量能双确认做空",
        "pros": ["快速参数适合白银高波动", "Force Index 过滤无量假信号", "OOS Return +56.62%"],
        "cons": ["快速参数容易过度交易", "夜盘流动性影响滑点", "做空天然面临无限风险"],
    },
    {
        "quadrant": "I Long",
        "instrument": "铁矿石 (I)",
        "direction": "做多",
        "regime": "Mild Trend",
        "freq": "1H",
        "version": "v23",
        "strategy_name": "mild_trend_long_I_1h_v23",
        "class_name": "MildTrendLongI1hV23",
        "research_path": "research/mild_trend/long/I/1h/v23_+32.21%",
        "indicators": [
            {"name": "Coppock Curve(10,14)", "category": "动量", "color": PURPLE,
             "desc": "Coppock 曲线 — 两个 ROC 的加权移动平均，从零线下方转正为买入信号"},
            {"name": "OI Flow(20)", "category": "持仓量", "color": GREEN,
             "desc": "OI Flow — 以价格方向加权的持仓量流向，OI Flow > Signal 表示机构在建仓"},
        ],
        "signal_logic": "Coppock > 0 且上升 + OI Flow > Signal → 动量反转+机构建仓确认",
        "pros": ["经典底部捕捉指标", "OI Flow 提供期货独有分析维度", "信号分层（强/弱）灵活"],
        "cons": ["Coppock 为滞后指标", "换月期间 OI 异常", "OOS Return +32.21% 偏保守"],
    },
    {
        "quadrant": "I Short",
        "instrument": "铁矿石 (I)",
        "direction": "做空",
        "regime": "Mild Trend",
        "freq": "2H",
        "version": "v3",
        "strategy_name": "mild_trend_short_I_2h_v3",
        "class_name": "MildTrendShortI2hV3",
        "research_path": "research/mild_trend/short/I/2h/v3_+21.19%",
        "indicators": [
            {"name": "HMA(40)", "category": "趋势", "color": GOLD,
             "desc": "Hull 移动均线 — 加权差值消除滞后，HMA 下降表示短期趋势转空"},
            {"name": "Schaff Trend Cycle(50,30)", "category": "动量", "color": PURPLE,
             "desc": "Schaff — MACD 经双重随机平滑，< 75 且下降表示动量衰竭"},
            {"name": "CMF(45)", "category": "成交量", "color": GREEN,
             "desc": "Chaikin Money Flow — CMF < 0 表示资金流出，卖方主导"},
        ],
        "signal_logic": "HMA 下降 + Schaff < 75 且下降 + CMF < 0 → 三重做空确认",
        "pros": ["HMA 低延迟入场及时", "三维信号过滤效果好（Sharpe 1.198）", "适合温和下跌"],
        "cons": ["中长周期参数可能错过日内反转", "Schaff 双重平滑反应偏慢", "政策驱动上涨时回撤较大"],
    },
]


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_indicator_card(slide, left, top, indicator):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(4.2), Inches(0.85),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()

    add_text_box(slide, left + 0.15, top + 0.05, 0.8, 0.25,
                 indicator["category"], font_size=10, color=indicator["color"], bold=True)
    add_text_box(slide, left + 1.0, top + 0.02, 3.0, 0.25,
                 indicator["name"], font_size=12, color=WHITE, bold=True)
    add_text_box(slide, left + 0.15, top + 0.3, 3.9, 0.55,
                 indicator["desc"], font_size=9, color=GRAY)


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, 0.5, 1.5, 9.0, 1.0,
                 "QBase_v2 策略研报", font_size=36, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 2.5, 9.0, 0.5,
                 "四象限最优策略分析", font_size=20, color=GOLD,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 3.3, 9.0, 0.5,
                 "AG Long  ·  AG Short  ·  I Long  ·  I Short", font_size=16, color=GRAY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 5.0, 9.0, 0.8,
                 "AlphaForge V7.2 Industrial Mode  |  2026-04-03",
                 font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


def create_overview_slide(prs, strategies):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, 0.5, 0.3, 9.0, 0.5,
                 "四象限策略总览", font_size=24, color=WHITE, bold=True)

    rows, cols = 5, 7
    tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.0), Inches(9.4), Inches(2.5)).table

    headers = ["象限", "品种/方向", "频率", "版本", "OOS Sharpe", "OOS Return", "指标组合"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = GOLD
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    for r, s in enumerate(strategies, 1):
        metrics = extract_metrics_from_html(PROJECT / s["research_path"] / "oos.html")
        sharpe = metrics.get("sharpe", "—")
        ret = metrics.get("return", "—")
        ind_names = " + ".join(ind["name"] for ind in s["indicators"])

        data = [s["quadrant"], f"{s['instrument']} {s['direction']}",
                s["freq"], s["version"], sharpe, f"+{ret}%" if not ret.startswith("-") else f"{ret}%", ind_names]
        for c, val in enumerate(data):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
                if c == 5:
                    p.font.color.rgb = GREEN if "+" in val else RED
                    p.font.bold = True

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if r > 0 else RGBColor(0x30, 0x30, 0x50)


def create_strategy_slides(prs, s):
    """Create 3 slides per strategy: metrics screenshot, equity screenshot, indicator analysis."""
    metrics = extract_metrics_from_html(PROJECT / s["research_path"] / "oos.html")

    # ── Slide 1: Metrics Screenshot ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide1, BG_DARK)

    add_text_box(slide1, 0.3, 0.15, 6.0, 0.4,
                 f"{s['quadrant']} — {s['instrument']} {s['direction']} | {s['freq']} {s['version']}",
                 font_size=18, color=WHITE, bold=True)
    add_text_box(slide1, 7.0, 0.15, 2.7, 0.4,
                 f"OOS Metrics", font_size=14, color=GOLD, alignment=PP_ALIGN.RIGHT)

    # Add metrics screenshot
    if s.get("metrics_img") and s["metrics_img"].exists():
        slide1.shapes.add_picture(str(s["metrics_img"]),
                                  Inches(0.2), Inches(0.6), Inches(9.6), Inches(6.2))
    else:
        add_text_box(slide1, 1.0, 3.0, 8.0, 1.0,
                     f"[报告截图: {s['research_path']}/oos.html]",
                     font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Equity Curve Screenshot ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, BG_DARK)

    add_text_box(slide2, 0.3, 0.15, 6.0, 0.4,
                 f"{s['quadrant']} — 权益曲线",
                 font_size=18, color=WHITE, bold=True)
    sharpe_str = metrics.get("sharpe", "—")
    ret_str = metrics.get("return", "—")
    add_text_box(slide2, 6.0, 0.15, 3.7, 0.4,
                 f"Sharpe {sharpe_str}  |  Return +{ret_str}%",
                 font_size=12, color=GREEN, alignment=PP_ALIGN.RIGHT)

    if s.get("equity_img") and s["equity_img"].exists():
        slide2.shapes.add_picture(str(s["equity_img"]),
                                  Inches(0.2), Inches(0.6), Inches(9.6), Inches(6.2))
    else:
        add_text_box(slide2, 1.0, 3.0, 8.0, 1.0,
                     "[权益曲线截图]", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 3: Indicator Analysis + Pros/Cons ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide3, BG_DARK)

    add_text_box(slide3, 0.3, 0.15, 9.0, 0.4,
                 f"{s['quadrant']} — 指标分析与优缺点",
                 font_size=18, color=WHITE, bold=True)

    # Signal logic
    add_text_box(slide3, 0.3, 0.6, 9.0, 0.25, "信号逻辑", font_size=12, color=GOLD, bold=True)
    add_text_box(slide3, 0.3, 0.85, 9.0, 0.35, s["signal_logic"], font_size=11, color=WHITE)

    # Indicators
    add_text_box(slide3, 0.3, 1.3, 9.0, 0.25, "指标组合", font_size=12, color=GOLD, bold=True)
    for i, ind in enumerate(s["indicators"]):
        add_indicator_card(slide3, 0.3 + (i % 2) * 4.6, 1.6 + (i // 2) * 0.95, ind)

    # Pros & Cons
    y_start = 1.6 + ((len(s["indicators"]) + 1) // 2) * 0.95 + 0.2

    add_text_box(slide3, 0.3, y_start, 4.5, 0.25, "✓ 优势", font_size=12, color=GREEN, bold=True)
    for i, pro in enumerate(s["pros"]):
        add_text_box(slide3, 0.5, y_start + 0.3 + i * 0.35, 4.2, 0.35,
                     f"· {pro}", font_size=9, color=WHITE)

    add_text_box(slide3, 5.2, y_start, 4.5, 0.25, "✗ 劣势", font_size=12, color=RED, bold=True)
    for i, con in enumerate(s["cons"]):
        add_text_box(slide3, 5.4, y_start + 0.3 + i * 0.35, 4.2, 0.35,
                     f"· {con}", font_size=9, color=WHITE)

    # Details
    y_details = y_start + 0.3 + max(len(s["pros"]), len(s["cons"])) * 0.35 + 0.2
    details = f"策略: {s['strategy_name']}  |  类: {s['class_name']}  |  信号: {'[0,1] 做多' if s['direction']=='做多' else '[-1,0] 做空'}  |  回测: Industrial Mode"
    add_text_box(slide3, 0.3, min(y_details, 6.8), 9.0, 0.3, details, font_size=8, color=GRAY)


def main():
    print("Taking screenshots...")
    opts = Options()
    opts.add_argument("--headless=new")
    opts.add_argument("--window-size=1400,900")
    opts.add_argument("--no-sandbox")
    opts.add_argument("--force-device-scale-factor=2")
    opts.binary_location = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

    driver = webdriver.Chrome(options=opts)

    for s in STRATEGIES:
        oos_html = PROJECT / s["research_path"] / "oos.html"
        if not oos_html.exists():
            print(f"  [skip] {s['quadrant']}: {oos_html}")
            continue

        # URL encode the path (handle + and % in folder name)
        url = "file://" + str(oos_html)
        print(f"  {s['quadrant']} {s['version']}...")
        driver.get(url)
        time.sleep(3)

        # Metrics screenshot (top of page)
        driver.set_window_size(1400, 900)
        driver.execute_script("window.scrollTo(0, 0)")
        time.sleep(1)
        metrics_path = SCREENSHOT_DIR / f"{s['strategy_name']}_metrics.png"
        driver.save_screenshot(str(metrics_path))
        s["metrics_img"] = metrics_path
        print(f"    metrics: {metrics_path.name}")

        # Equity curve screenshot (scroll to section)
        driver.execute_script("window.scrollTo(0, 680)")
        time.sleep(1)
        equity_path = SCREENSHOT_DIR / f"{s['strategy_name']}_equity.png"
        driver.save_screenshot(str(equity_path))
        s["equity_img"] = equity_path
        print(f"    equity: {equity_path.name}")

    driver.quit()

    print("\nBuilding PPT...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_overview_slide(prs, STRATEGIES)

    for s in STRATEGIES:
        create_strategy_slides(prs, s)

    output = Path("/Users/simon/Desktop/QBase_Strategy_Report.pptx")
    prs.save(str(output))
    print(f"\nPPT saved: {output}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
